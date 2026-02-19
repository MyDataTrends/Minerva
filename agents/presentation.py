"""
Presentation Agent — The "Closer".

Converts Markdown Sales Kits into polished PowerPoint presentations.
"""

from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from agents.base import AgentConfig, AgentResult, BaseAgent, TriggerType
from agents.memory.operational import OperationalMemory
from utils.logging import get_logger

logger = get_logger(__name__)

# Directory where sales kits are stored
SALES_KIT_DIR = Path(__file__).resolve().parents[1] / "agents" / "knowledge_base" / "product" / "sales_kits"
# Directory to save presentations
PRESENTATION_DIR = Path(__file__).resolve().parents[1] / "reports" / "presentations"


class PresentationAgent(BaseAgent):
    """
    Automated Slide Deck Generator.
    
    Workflow:
    1. Finds the latest (or specified) Markdown Sales Kit.
    2. Parses the Markdown structure (Headers -> Slide Titles, Bullets -> Content).
    3. Generates a .pptx file using a clean, professional template.
    4. Saves to reports/presentations.
    """

    name = "presentation"
    trigger_type = TriggerType.MANUAL

    def __init__(self, config: Optional[AgentConfig] = None):
        super().__init__(config)
        self.memory = OperationalMemory("presentation")
        PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)

    def run(self, **kwargs) -> AgentResult:
        """Execute the presentation generation workflow."""
        result = self._make_result()
        
        # 1. Find Input File
        target_file = kwargs.get("file")
        if not target_file:
            target_file = self._find_latest_sales_kit()
            
        if not target_file:
            result.error = "No Sales Kit Markdown files found."
            result.success = False
            return result
        
        # Resolve to Path object if it's a string
        if isinstance(target_file, str):
            target_file = Path(target_file)
            
        result.add_action(f"selected_kit:{target_file.name}")
        
        # 2. Parse Markdown
        try:
            content = target_file.read_text(encoding="utf-8")
            slides_data = self._parse_markdown_to_slides(content)
            result.add_action(f"parsed_{len(slides_data)}_slides")
        except Exception as e:
            result.error = f"Failed to parse markdown: {e}"
            result.success = False
            return result

        # 3. Generate PPTX
        try:
            output_path = self._generate_pptx(slides_data, target_file.stem)
            result.add_action(f"generated_pptx:{output_path.name}")
            result.summary = f"Generated Slide Deck: {output_path.name} ({len(slides_data)} slides)"
            
            result.metrics["slides_generated"] = len(slides_data)
            result.metrics["output_file"] = str(output_path)
            
        except Exception as e:
            result.error = f"Failed to generate PPTX: {e}"
            # Check for missing dependency
            if "No module named 'pptx'" in str(e):
                result.error = "Missing dependency: python-pptx. Please install it."
            result.success = False
            return result

        return result

    def _find_latest_sales_kit(self) -> Optional[Path]:
        """Find the most recently created Sales Kit markdown file."""
        if not SALES_KIT_DIR.exists():
            return None
        candidates = list(SALES_KIT_DIR.glob("*.md"))
        if not candidates:
            return None
        candidates.sort(key=lambda p: p.stat().st_mtime, reverse=True)
        return candidates[0]

    def _parse_markdown_to_slides(self, markdown_text: str) -> List[dict]:
        """Parse Markdown into a list of slide dictionaries."""
        slides = []
        lines = markdown_text.split('\n')
        
        current_slide = None
        
        # Helper to clean text
        def clean(text):
            # Remove bold/italic markers
            text = re.sub(r'[\*_]{2}(.*?)[\*_]{2}', r'\1', text) # **bold**
            text = re.sub(r'[\*_](.*?)[\*_]', r'\1', text)       # *italic*
            return text.strip()
        
        # Extract Main Title from first H1
        first_h1_match = re.search(r'^#\s+(.+)$', markdown_text, re.MULTILINE)
        if first_h1_match:
            main_title = clean(first_h1_match.group(1))
            subtitle = ""
            for line in lines:
                if line.startswith(">"):
                    subtitle += line.replace(">", "").strip() + "\n"
            
            slides.append({
                "type": "title", 
                "title": main_title, 
                "subtitle": clean(subtitle)
            })

        # Parse content slides
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # H2 starts a new slide
            if line.startswith("## "):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "type": "content",
                    "title": clean(line.replace("## ", "")),
                    "content": []
                }
            
            # H3 can also start a new slide or be a sub-header
            elif line.startswith("### "):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "type": "content",
                    "title": clean(line.replace("### ", "")),
                    "content": []
                }
                
            # Bullets
            elif line.startswith(("* ", "- ")):
                if current_slide:
                    clean_line = clean(re.sub(r'^[\*\-]\s+', '', line))
                    current_slide["content"].append(clean_line)
                    
            # Text block (if it looks like content)
            elif current_slide and not line.startswith(("#", ">", "---")):
                 current_slide["content"].append(clean(line))

        # Append last slide
        if current_slide:
            slides.append(current_slide)
            
        return slides

    def _generate_pptx(self, slides_data: List[dict], filename_base: str) -> Path:
        """Create the PowerPoint file with styling."""
        prs = Presentation()
        
        # Define Assay Colors
        DARK_BG = (20, 24, 35)      # Deep Dark Blue/Black
        ACCENT_COLOR = (0, 190, 255) # Cyan/Blue Highlight
        TEXT_COLOR = (240, 240, 240) # Off-white
        
        from pptx.dml.color import RGBColor
        
        def set_background(slide):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*DARK_BG)
            
        def style_title(shape, is_main=False):
            if not shape: return
            shape.text_frame.word_wrap = True
            p = shape.text_frame.paragraphs[0]
            p.font.color.rgb = RGBColor(*ACCENT_COLOR) if is_main else RGBColor(*TEXT_COLOR)
            p.font.bold = True
            p.font.name = "Arial"
            p.font.size = Pt(44) if is_main else Pt(32)

        def style_subtitle(shape):
            if not shape: return
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(*TEXT_COLOR)
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(24)

        def style_body(shape):
            if not shape: return
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(*TEXT_COLOR)
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(18)
                paragraph.space_after = Pt(10)

        # 1. Main Title Slide
        title_slide_layout = prs.slide_layouts[0]
        bullet_slide_layout = prs.slide_layouts[1]
        
        for slide_info in slides_data:
            if slide_info["type"] == "title":
                slide = prs.slides.add_slide(title_slide_layout)
                set_background(slide)
                
                title = slide.shapes.title
                title.text = slide_info.get("title", "Presentation")
                style_title(title, is_main=True)
                
                subtitle = slide.placeholders[1]
                subtitle.text = slide_info.get("subtitle", "")
                style_subtitle(subtitle)
                
            elif slide_info["type"] == "content":
                slide = prs.slides.add_slide(bullet_slide_layout)
                set_background(slide)
                
                shapes = slide.shapes
                title_shape = shapes.title
                title_shape.text = slide_info.get("title", "")
                style_title(title_shape, is_main=False)
                # Adjust title color for content slides to be accent? Or keep white?
                # Let's make content titles cyan for pizzazz
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*ACCENT_COLOR)
                
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.word_wrap = True
                
                content_items = slide_info.get("content", [])
                
                # Clear existing paragraphs (often one empty one exists)
                tf.clear() 

                for item in content_items:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
                
                style_body(body_shape)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = PRESENTATION_DIR / f"{filename_base}_{timestamp}.pptx"
        prs.save(output_file)
        
        return output_file
